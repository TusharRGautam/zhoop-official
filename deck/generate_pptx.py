"""
Python PPTX Generator for Digital Solutions Investor Pitch Deck (10 Slides)
Generates an editable 16:9 widescreen presentation with custom shapes, tables, and image placeholders.
"""

import sys
import subprocess

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

def create_pitch_deck(filename="Digital_Solutions_Pitch_Deck.pptx"):
    prs = Presentation()
    
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Color Palette
    COLOR_PRIMARY = RGBColor(79, 70, 229)    # #4f46e5 Indigo
    COLOR_DARK = RGBColor(15, 23, 42)        # #0f172a Dark text
    COLOR_MUTED = RGBColor(71, 85, 105)      # #475569 Muted text
    COLOR_BG_CARD = RGBColor(248, 250, 252)  # #f8fafc Light card
    COLOR_EMERALD = RGBColor(16, 185, 129)   # #10b981 Emerald
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_BORDER = RGBColor(226, 232, 240)
    
    def add_header(slide, slide_num, category, title):
        # Category Tag
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        tf = tag_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{slide_num.upper()}  •  {category.upper()}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK

    def add_card(slide, left, top, width, height, bg_color=COLOR_BG_CARD, border_color=COLOR_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    def add_placeholder_box(slide, left, top, width, height, text="[ Image / Screenshot Placeholder ]"):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(241, 245, 249)
        shape.line.color.rgb = RGBColor(203, 213, 225)
        shape.line.width = Pt(1.5)
        
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_MUTED
        return shape

    # -------------------------------------------------------------
    # SLIDE 1: COVER
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    
    # Left Hero Container
    left_card = add_card(s1, Inches(0.8), Inches(0.8), Inches(7.2), Inches(5.9), COLOR_WHITE, COLOR_BORDER)
    tf = left_card.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "EXECUTIVE PITCH DECK  •  SEED STAGE"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    p2 = tf.add_paragraph()
    p2.text = "Zhoop"
    p2.font.size = Pt(38)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_DARK
    
    p3 = tf.add_paragraph()
    p3.text = "Building Digital Solutions for Growing Businesses"
    p3.font.size = Pt(18)
    p3.font.color.rgb = COLOR_MUTED
    
    p4 = tf.add_paragraph()
    p4.text = "\nWebsites  •  CRM  •  Apps  •  E-commerce  •  HTML5 Games  •  Automation\n"
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_PRIMARY
    
    p5 = tf.add_paragraph()
    p5.text = "₹15 Lakh Investment Ask\nHigh-efficiency capital round focused on sales engine & market expansion."
    p5.font.size = Pt(14)
    p5.font.bold = True
    p5.font.color.rgb = COLOR_EMERALD

    # Right Hero Image Placeholder
    add_placeholder_box(s1, Inches(8.3), Inches(0.8), Inches(4.2), Inches(5.9), "[ Cover Image / Product Mockup Placeholder ]\n\n(Right-click -> Change Picture to insert image)")

    # -------------------------------------------------------------
    # SLIDE 2: PROBLEM
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "Slide 02", "The Market Problem", "Businesses Need Digital Products, But Development Is Fragmented")
    
    problems = [
        ("Custom Dev is Expensive", "Agencies quote exorbitant fees and lengthy development cycles out of reach for SMEs."),
        ("Generic SaaS Doesn't Fit", "Rigid off-the-shelf software doesn't adapt to unique business operations."),
        ("Multiple Vendors Required", "Managing separate vendors for Web, CRM, and Apps causes friction and delays."),
        ("Development is Slow", "Lack of reusable components leads to extended delivery turnarounds."),
        ("Lack of Reliable Partners", "SMEs suffer from lack of post-launch maintenance and freelancer unreliability."),
        ("Missing Engaging Experiences", "Businesses miss out on interactive gamified marketing and HTML5 brand games.")
    ]
    
    col_w, row_h = Inches(3.7), Inches(2.0)
    start_x, start_y = Inches(0.8), Inches(1.8)
    
    for i, (title, desc) in enumerate(problems):
        r, c = i // 3, i % 3
        card = add_card(s2, start_x + c * (col_w + Inches(0.25)), start_y + r * (row_h + Inches(0.2)), col_w, row_h)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_MUTED

    # Gap Box
    gap_box = add_card(s2, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.8), RGBColor(238, 242, 255), COLOR_PRIMARY)
    tf = gap_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THE GAP: One affordable technology partner that can build different types of digital products under one roof."
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # -------------------------------------------------------------
    # SLIDE 3: SOLUTION
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Slide 03", "Our Value Proposition", "One Technology Partner for All Digital Products")
    
    solutions = [
        ("Websites", "Corporate websites, high-converting landing pages, and interactive client portals."),
        ("Custom CRM", "Lead management, sales pipeline, business operations, and automated dashboards."),
        ("Mobile Apps", "Customer, staff, booking, and on-demand service mobile applications."),
        ("HTML5 Games", "Browser-based games, branded viral games, and promotional interactive campaigns."),
        ("E-commerce", "Online stores, payment gateways, inventory sync & multi-channel API integrations."),
        ("Automation", "WhatsApp bots, API connectors, and automated business workflow pipelines.")
    ]
    
    for i, (title, desc) in enumerate(solutions):
        r, c = i // 3, i % 3
        card = add_card(s3, start_x + c * (col_w + Inches(0.25)), start_y + r * (row_h + Inches(0.2)), col_w, row_h)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_MUTED

    pipe_box = add_card(s3, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.8), COLOR_BG_CARD)
    tf = pipe_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Unified Delivery:  Website  →  CRM  →  App  →  Game  →  Automation  →  Long-Term Support"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK

    # -------------------------------------------------------------
    # SLIDE 4: BUSINESS MODEL
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "Slide 04", "Monetization & Pricing", "Multiple Revenue Streams with High Expansion Value")
    
    # Table on Left
    rows, cols = 8, 3
    t_shape = s4.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8))
    table = t_shape.table
    table.columns[0].width = Inches(2.6)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(1.7)
    
    headers = ["Service Vertical", "Typical Opportunity", "Type"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(241, 245, 249)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK

    data = [
        ("Websites", "₹20K – ₹2L+", "Land / Entry"),
        ("Custom CRM", "₹50K – ₹5L+", "High Value"),
        ("Mobile Apps", "₹1L – ₹10L+", "Enterprise"),
        ("E-commerce", "₹50K – ₹5L+", "Growth"),
        ("HTML5 Games", "₹25K – ₹5L+", "High Margin"),
        ("Automation", "Additional Revenue", "Upsell"),
        ("Maintenance & Support", "Recurring Retainer", "Recurring")
    ]
    
    for r_idx, row in enumerate(data, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            if c_idx == 1:
                p.font.bold = True
                p.font.color.rgb = COLOR_PRIMARY

    # Right Card: Land Expand Retain
    right_card = add_card(s4, Inches(7.6), Inches(1.8), Inches(4.9), Inches(4.8), COLOR_WHITE, COLOR_BORDER)
    tf = right_card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "CUSTOMER EXPANSION ENGINE\n"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    p2 = tf.add_paragraph()
    p2.text = "1. LAND\nStart with one project (Website or Automation).\n"
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_DARK
    
    p3 = tf.add_paragraph()
    p3.text = "2. EXPAND\nUpsell Custom CRM, Mobile App, or HTML5 Games.\n"
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_DARK
    
    p4 = tf.add_paragraph()
    p4.text = "3. RETAIN\nRecurring monthly retainer for maintenance, hosting & support.\n\n"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_DARK
    
    p5 = tf.add_paragraph()
    p5.text = "💡 'A single client acquired can generate multiple projects across their business journey.'"
    p5.font.size = Pt(11)
    p5.font.italic = True
    p5.font.color.rgb = COLOR_MUTED

    # -------------------------------------------------------------
    # SLIDE 5: MARKET OPPORTUNITY
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "Slide 05", "Market Demand & Verticals", "Growing Demand for Digital Business Solutions")
    
    m_left = add_card(s5, Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.8), COLOR_WHITE, COLOR_BORDER)
    tf = m_left.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "TARGET CUSTOMER VERTICALS\n"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    industries = "• Real Estate\n• Automotive\n• Healthcare\n• Hospitality\n• Retail & E-commerce\n• Professional Services\n• Tech Startups\n• Local & Service Businesses\n• Brands & Marketing Agencies"
    p2 = tf.add_paragraph()
    p2.text = industries
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_DARK

    m_right = add_card(s5, Inches(7.1), Inches(1.8), Inches(5.4), Inches(4.8), RGBColor(250, 245, 255), RGBColor(233, 213, 255))
    tf = m_right.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "HIGH-MARGIN NICHE: HTML5 GAMES\n"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(124, 58, 237)
    
    p2 = tf.add_paragraph()
    p2.text = "HTML5 games unlock major enterprise marketing opportunities in:\n\n• Marketing & Festive Campaigns\n• Brand Engagement & Virality\n• Promotional Giveaways & Gamified Contests\n• Playable Advertising & High-Retention Activations\n\n[ Play Live HTML5 Games Arcade: zhoop.in/games.html ]"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_DARK

    # -------------------------------------------------------------
    # SLIDE 6: GO-TO-MARKET
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "Slide 06", "Customer Acquisition", "Founder-Led Today. Sales-Led Tomorrow.")
    
    gtm_1 = add_card(s6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.0), COLOR_BG_CARD)
    tf = gtm_1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "CURRENT STAGE: 1 FOUNDER + 1 CO-FOUNDER\n"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_MUTED
    
    p2 = tf.add_paragraph()
    p2.text = "• Sales + Tech Architecture + Management\n• Direct founder outreach & pitching\n• Establishing repeatable execution playbooks"
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_DARK

    gtm_2 = add_card(s6, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.0), RGBColor(238, 242, 255), COLOR_PRIMARY)
    tf = gtm_2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AFTER ₹15L INVESTMENT: 2 SALES EXECUTIVES\n"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    p2 = tf.add_paragraph()
    p2.text = "• Outbound Calling & Cold Prospecting\n• WhatsApp Marketing Funnels\n• LinkedIn Decision-Maker Outreach\n• Agency & Partner Referrals\n• Local Business Campaigns\n• Meta & Google High-Intent Ads"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_DARK

    gtm_bar = add_card(s6, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.8), COLOR_WHITE, COLOR_BORDER)
    tf = gtm_bar.text_frame
    p = tf.paragraphs[0]
    p.text = "Acquisition Goal: Lead  →  Client  →  Project  →  Repeat Business  |  Paid Ads + Outbound + Referrals"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    # -------------------------------------------------------------
    # SLIDE 7: SCALING STRATEGY
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Slide 07", "Capital-Efficient Scaling", "₹15L → Revenue → Developers → Higher Capacity")
    
    # 4 Phase Cards
    phases = [
        ("PHASE 1", "1 Founder + 1 Co-Founder", "Establish client proof & initial delivery playbooks."),
        ("PHASE 2", "2 Sales + 2 Creative", "Funded by ₹15L round to scale deal pipeline."),
        ("PHASE 3", "More Clients & Cashflow", "Accelerating revenue from increased project inflow."),
        ("PHASE 4", "2 Full-Stack Devs", "Hired sustainably from generated cashflow.")
    ]
    p_w = Inches(2.7)
    for i, (ph_num, title, desc) in enumerate(phases):
        card = add_card(s7, Inches(0.8) + i * (p_w + Inches(0.3)), Inches(1.8), p_w, Inches(2.3), COLOR_WHITE, COLOR_BORDER)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ph_num
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_DARK
        
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(10)
        p3.font.color.rgb = COLOR_MUTED

    dev_box = add_card(s7, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.2), RGBColor(238, 242, 255), COLOR_PRIMARY)
    tf = dev_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💡 REVENUE-FUNDED DEVELOPER HIRING (Not from ₹15L Seed)\n"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    p2 = tf.add_paragraph()
    p2.text = "• Developers initial ₹15L investment ka part nahi hain.\n• They will be hired through generated revenue and project reinvestment as volume grows.\n• Estimated Developer Cost: ₹30K–₹40K/month per developer (2 Devs: ₹7.2L–₹9.6L/year)."
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_DARK

    # -------------------------------------------------------------
    # SLIDE 8: INVESTMENT ASK
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "Slide 08", "Fund Allocation", "₹15 Lakh Investment Ask")
    
    t_shape = s8.shapes.add_table(7, 3, Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8))
    table = t_shape.table
    table.columns[0].width = Inches(3.0)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(1.5)
    
    for i, h in enumerate(["Use of Funds", "Amount", "% of Round"]):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(241, 245, 249)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK

    funds_data = [
        ("2 Sales Executives", "₹4.8 Lakh", "32.0%"),
        ("2 Creative / Designers", "₹4.8 Lakh", "32.0%"),
        ("Office Infrastructure", "₹2.4 Lakh", "16.0%"),
        ("Security Deposit", "₹1.0 Lakh", "6.7%"),
        ("Advertising & Meta Ads", "₹2.0 Lakh", "13.3%"),
        ("TOTAL INVESTMENT", "₹15.0 Lakh", "100%")
    ]
    for r_idx, row in enumerate(funds_data, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            if r_idx == 6 or c_idx == 1:
                p.font.bold = True
                p.font.color.rgb = COLOR_PRIMARY

    strat_box = add_card(s8, Inches(7.6), Inches(1.8), Inches(4.9), Inches(4.8), RGBColor(30, 27, 75))
    tf = strat_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "INVESTMENT STRATEGY\n"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(165, 180, 252)
    
    p2 = tf.add_paragraph()
    p2.text = "₹15L Seed Capital\n      ↓\nSales + Marketing Engine\n      ↓\nMore Client Projects\n      ↓\nRevenue Reinvestment → 2 Devs\n      ↓\nScalable High-Margin Delivery"
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE

    # -------------------------------------------------------------
    # SLIDE 9: WHY US
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "Slide 09", "Competitive Advantage", "Built for Capital-Efficient Growth")
    
    why_us = [
        ("01 — Lean Team", "No unnecessary early hiring. Lean operations ensure long runway."),
        ("02 — Sales First", "Investment is strictly focused on generating revenue and clients."),
        ("03 — Multiple Services", "Websites, CRM, Apps, HTML5 Games, E-commerce & Automation."),
        ("04 — Cross-Selling", "One customer generates multiple project purchases."),
        ("05 — Revenue-Funded Hiring", "Developers are added when verified cashflow justifies the cost."),
        ("06 — Reusable Tech", "Pre-built components and modular systems accelerate delivery.")
    ]
    
    for i, (title, desc) in enumerate(why_us):
        r, c = i // 3, i % 3
        card = add_card(s9, start_x + c * (col_w + Inches(0.25)), start_y + r * (row_h + Inches(0.2)), col_w, row_h)
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_MUTED

    # -------------------------------------------------------------
    # SLIDE 10: PROJECTIONS & VISION
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "Slide 10", "Financial Targets & Vision", "₹50L → ₹1Cr+ Annual Revenue Target")
    
    t_shape = s10.shapes.add_table(8, 3, Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8))
    table = t_shape.table
    table.columns[0].width = Inches(2.7)
    table.columns[1].width = Inches(1.9)
    table.columns[2].width = Inches(1.9)
    
    for i, h in enumerate(["Revenue Source", "₹50L Target", "₹1Cr Target"]):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(241, 245, 249)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK

    proj_data = [
        ("Websites", "₹10 Lakh", "₹20 Lakh"),
        ("CRM Software", "₹13 Lakh", "₹27 Lakh"),
        ("Mobile Apps", "₹10 Lakh", "₹20 Lakh"),
        ("HTML5 Games", "₹7 Lakh", "₹15 Lakh"),
        ("E-commerce & Automation", "₹6 Lakh", "₹10 Lakh"),
        ("Maintenance & Recurring", "₹4 Lakh", "₹8 Lakh"),
        ("TOTAL REVENUE", "₹50 Lakh", "₹1.00 Crore")
    ]
    for r_idx, row in enumerate(proj_data, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            if r_idx == 7 or c_idx >= 1:
                p.font.bold = True
                if r_idx == 7 and c_idx == 2:
                    p.font.color.rgb = COLOR_EMERALD
                elif c_idx >= 1:
                    p.font.color.rgb = COLOR_PRIMARY

    v_box = add_card(s10, Inches(7.6), Inches(1.8), Inches(4.9), Inches(4.8), COLOR_WHITE, COLOR_BORDER)
    tf = v_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "OUR VISION\n"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    p2 = tf.add_paragraph()
    p2.text = "Build a scalable technology company delivering websites, software, apps and interactive digital experiences for growing businesses.\n\n"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_DARK
    
    p3 = tf.add_paragraph()
    p3.text = "⚡ GROWTH FLYWHEEL:\n₹15L Seed → Sales Team → Client Volume → Reinvested Developers → Higher Delivery → ₹1Cr+ Target"
    p3.font.size = Pt(11)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_EMERALD

    prs.save(filename)
    print(f"[SUCCESS] Successfully created '{filename}' with all 10 slides!")

if __name__ == "__main__":
    create_pitch_deck()
